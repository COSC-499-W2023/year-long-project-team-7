from enum import Enum
from io import BytesIO
import json
import typing
from pptx import Presentation  # type: ignore
from django.core.files.storage import FileSystemStorage
from django.core.files.base import ContentFile
import random
from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE
from .models import File
from .utils import error
from typing import List, Optional, Dict, Any


class FieldTypes(Enum):
    TITLE = "TITLE"
    TEXT = "TEXT"
    IMAGE = "IMAGE"


class SlideTypes(Enum):
    TITLE = "TITLE"
    CONTENT = "CONTENT"
    IMAGE = "IMAGE"
    MULTIPLE_CHOICE = "MULTIPLE_CHOICE"
    TRUE_FALSE = "TRUE_FALSE"
    SHORT_ANSWER = "SHORT_ANSWER"


class SlideField:
    def __init__(self, field_index: int, field_type: FieldTypes, value: str):
        self.field_index = field_index
        self.field_type = field_type
        self.value = value


class SlideContent:
    def __init__(
        self,
        slide_num: Optional[int] = None,
        layout: Optional[str] = None,
        fields: Optional[List[SlideField]] = None,
        json: Optional[Dict[Any, Any]] = None,
    ):
        if json is not None:
            self.slide_num = json["SLIDE_NUM"]
            self.layout = json["SLIDE_LAYOUT"]
            self.fields = [
                SlideField(
                    field["FIELD_INDEX"],
                    FieldTypes(field["FIELD_TYPE"]),
                    field["FIELD_VALUE"],
                )
                for field in json["FIELDS"]
            ]
        else:
            self.slide_num = slide_num or 0
            self.layout = layout or ""
            self.fields = fields or []

    def to_json_string(self) -> str:
        slide_json = {
            "SLIDE_NUM": self.slide_num,
            "SLIDE_LAYOUT": self.layout,
            "FIELDS": [
                {
                    "FIELD_INDEX": field.field_index,
                    "FIELD_TYPE": field.field_type.value,
                    "FIELD_VALUE": field.value,
                }
                for field in self.fields
            ],
        }
        return str(slide_json)

    def to_json(self) -> dict[str, str]:
        slide_json = {
            "SLIDE_NUM": self.slide_num,
            "SLIDE_LAYOUT": self.layout,
            "FIELDS": [
                {
                    "FIELD_INDEX": field.field_index,
                    "FIELD_TYPE": field.field_type.value,
                    "FIELD_VALUE": field.value,
                }
                for field in self.fields
            ],
        }
        return slide_json

    def update_from_json(self, response: str) -> None:
        object = json.loads(response)
        for field in self.fields:
            for new_field in object["FIELDS"]:
                if str(field.field_index) == new_field["FIELD_INDEX"]:
                    field.value = new_field["FIELD_VALUE"]


class MissingPlaceholderError(Exception):
    def __init__(self, message: str):
        self.message = "Error: " + message
        super().__init__(self.message)


PLACEHOLDERS = {
    SlideTypes.TITLE: {
        FieldTypes.TITLE: "<PRESENTATION TITLE HERE>",
        FieldTypes.TEXT: "<PRESENTATION SUBTITLE HERE>",
    },
    SlideTypes.CONTENT: {
        FieldTypes.TITLE: "<SLIDE TITLE HERE>",
        FieldTypes.TEXT: "<SLIDE TEXT HERE>",
    },
    SlideTypes.IMAGE: {
        FieldTypes.TITLE: "<SLIDE TITLE HERE>",
        FieldTypes.TEXT: "<SLIDE TEXT HERE>",
        FieldTypes.IMAGE: "<IMAGE SEARCH QUERY HERE>",
    },
    SlideTypes.MULTIPLE_CHOICE: {
        FieldTypes.TITLE: "<SLIDE QUESTION HERE>",
        FieldTypes.TEXT: "A) <SLIDE ANSWER HERE>  B) <SLIDE ANSWER HERE>  C) <SLIDE ANSWER HERE>  D) <SLIDE ANSWER HERE>",
    },
    SlideTypes.TRUE_FALSE: {
        FieldTypes.TITLE: "<SLIDE QUESTION HERE>",
        FieldTypes.TEXT: "A) TRUE  B) FALSE",
    },
    SlideTypes.SHORT_ANSWER: {
        FieldTypes.TITLE: "<SLIDE QUESTION HERE>",
        FieldTypes.TEXT: "<SLIDE ANSWER HERE>",
    },
}


class PresentationManager:
    def __init__(self, template: File) -> None:
        file_system = FileSystemStorage()

        self.template_path = file_system.path(template.file.name)
        self.presentation = Presentation(self.template_path)

    def delete_all_slides(self) -> None:
        # Delete all slides from template presentation
        xml_slides = self.presentation.slides._sldIdLst
        slides = list(xml_slides)
        for slide in slides:
            rId = slide.get("r:id")
            if rId is not None:
                self.presentation.part.drop_rel(rId)
        xml_slides.clear()

    def save_presentation(self, conversion_id: int) -> str:
        file_system = FileSystemStorage()
        buffer = BytesIO()
        self.presentation.save(buffer)
        buffer.seek(0)
        file_content = ContentFile(buffer.read())

        rel_path = f"conversion_output_{conversion_id}.pptx"

        file_system.save(rel_path, file_content)

        return rel_path

    def get_title_slide_layout(self) -> typing.Any:
        potential_layouts = []

        temp_presentation = Presentation(self.template_path)

        for slide_layout in temp_presentation.slide_layouts:
            if "title" in slide_layout.name.lower():
                potential_layouts.append(slide_layout)
        try:
            return random.choice(potential_layouts)
        except IndexError:
            raise MissingPlaceholderError("Template does not have a title slide")

    def get_image_slide_layout(self) -> typing.Any:
        potential_layouts = []

        temp_presentation = Presentation(self.template_path)

        for slide_layout in temp_presentation.slide_layouts:
            temp_slide = temp_presentation.slides.add_slide(slide_layout)

            for shape in temp_slide.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == PP_PLACEHOLDER.PICTURE:
                        potential_layouts.append(slide_layout)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    potential_layouts.append(slide_layout)
        try:
            return random.choice(potential_layouts)
        except IndexError:
            raise MissingPlaceholderError(
                "Template does not have any image placeholders"
            )

    def get_content_slide_layout(self) -> typing.Any:
        potential_layouts = []

        temp_presentation = Presentation(self.template_path)

        for slide_layout in temp_presentation.slide_layouts:
            temp_slide = temp_presentation.slides.add_slide(slide_layout)

            for shape in temp_slide.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == PP_PLACEHOLDER.BODY:
                        potential_layouts.append(slide_layout)
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    potential_layouts.append(slide_layout)
        try:
            return random.choice(potential_layouts)
        except IndexError:
            raise MissingPlaceholderError(
                "Template does not have any text placeholders"
            )

    def get_field(
        self, field_index: int, slide_type: SlideTypes, field_type: FieldTypes
    ) -> SlideField:
        default = PLACEHOLDERS.get(slide_type, {})
        placeholder = default.get(field_type, f"<{field_type.value} HERE>")
        return SlideField(field_index, field_type, placeholder)

    def get_slide_layout_fields(
        self, layout: typing.Any, slide_type: SlideTypes
    ) -> list[SlideField]:
        fields = []
        temp_presentation = Presentation(self.template_path)
        temp_slide = temp_presentation.slides.add_slide(layout)
        shapes = temp_slide.shapes
        for shape in shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format

                if phf.type == PP_PLACEHOLDER.TITLE:
                    fields.append(
                        self.get_field(
                            shapes.index(shape), slide_type, FieldTypes.TITLE
                        )
                    )

                elif phf.type == PP_PLACEHOLDER.BODY:
                    fields.append(
                        self.get_field(shapes.index(shape), slide_type, FieldTypes.TEXT)
                    )

                elif phf.type == PP_PLACEHOLDER.PICTURE:
                    fields.append(
                        self.get_field(
                            shapes.index(shape), slide_type, FieldTypes.IMAGE
                        )
                    )
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                fields.append(
                    self.get_field(shapes.index(shape), slide_type, FieldTypes.TEXT)
                )

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                fields.append(
                    self.get_field(shapes.index(shape), slide_type, FieldTypes.IMAGE)
                )

        return fields

    def add_slide_to_presentation(self, slide_content: SlideContent) -> None:
        layout = self.presentation.slide_layouts.get_by_name(slide_content.layout)
        slide = self.presentation.slides.add_slide(layout)
        fields = slide_content.fields

        for field in fields:
            if field.field_type == FieldTypes.TITLE:
                slide.shapes[field.field_index].text = field.value

            elif field.field_type == FieldTypes.TEXT:
                slide.shapes[field.field_index].text = field.value

            elif field.field_type == FieldTypes.IMAGE:
                try:
                    slide.shapes[field.field_index].insert_picture(field.value)
                except Exception as e:
                    # ValueError: unsupported image format, expected one of: dict_keys(['BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF']), got 'WEBP'
                    # Other random errors
                    error(e)
                    continue
