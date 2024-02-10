import typing
from pptx import Presentation  # type: ignore

from .presentationManager import (
    PresentationManager,
    SlideFieldTypes,
    SlideContent,
    SlideTypes,
)
from .models import Conversion
from django.conf import settings
import json
import os
from django.core.files.storage import FileSystemStorage
from django.core.files.base import ContentFile
from io import BytesIO
import re
import random
from openai import OpenAI
import requests
from .models import File
from concurrent.futures import ThreadPoolExecutor
from tenacity import retry, stop_after_attempt, wait_random_exponential
from enum import Enum


class PresentationGenerator:
    def __init__(self, input_file_paths: list[str], conversion: Conversion):
        self.manager = PresentationManager()

        with open(os.path.join(settings.BASE_DIR, "prompts.json"), "r") as file:
            self.prompts = json.load(file)

        self.conversion = conversion
        self.user_prompt = json.loads(conversion.user_parameters).get("prompt", "")
        self.language = json.loads(conversion.user_parameters).get("language", "")
        self.tone = json.loads(conversion.user_parameters).get("tone", "")
        self.complexity = json.loads(conversion.user_parameters).get("complexity", 3)
        self.num_slides = json.loads(conversion.user_parameters).get("num_slides", 10)
        self.image_frequency = json.loads(conversion.user_parameters).get(
            "image_frequency", 3
        )
        self.template = json.loads(conversion.user_parameters).get("template", 1)

        self.client = OpenAI(api_key=settings.OPENAI_API_KEY, organization="org-hYPJO2WKqFAN1o75AdfdfeBx")  # type: ignore

        self.openai_files = []
        for path in input_file_paths:
            self.openai_files.append(
                self.client.files.create(file=open(path, "rb"), purpose="assistants")
            )

        instructions = f"{self.prompts['reader']}"

        instructions = (
            instructions
            if self.language == "Auto"
            else f"{instructions} Only respond in {self.language}"
        )

        instructions = (
            instructions
            if self.tone == "Auto"
            else f"{instructions} {self.prompts['tone'].format(tone=self.tone)}"
        )

        instructions = (
            instructions
            if self.complexity == 3
            else f"{instructions} {self.prompts['complexity'].format(complexity=self.complexity)}"
        )

        instructions = (
            instructions
            if len(self.user_prompt) == 0 or self.user_prompt == ""
            else f"{instructions} {self.prompts['prompt-input'].format(prompt=self.user_prompt)}"
        )

        self.assistant = self.client.beta.assistants.create(
            instructions=instructions,
            tools=[{"type": "retrieval"}],
            model=settings.OPENAI_MODEL,  # type: ignore
            file_ids=[file.id for file in self.openai_files],
        )

    def prompt_assistant(self, proompt: str) -> str:
        thread = self.create_thread()
        self.add_message_to_thread(thread.id, proompt)
        run = self.create_run(thread.id)

        while run.status != "completed":
            run = self.retrieve_run(thread.id, run.id)
            if run.status == "failed":
                if run.last_error:
                    print(run.last_error.message)
                    return str(run.last_error.message)

        messages = []

        for msg in self.list_messages(thread.id):
            if msg.role == "assistant":
                content = msg.content[0]
                if hasattr(content, "text"):
                    messages.append(
                        re.sub(r"【.*】", "", content.text.value)
                    )  # Remove source links

        self.delete_thread(thread.id)
        return "\n".join(messages)

    def image_search(self, query: str) -> File:
        api_url = "https://api.unsplash.com/search/photos"
        params = {"query": query, "client_id": settings.UNSPLASH_ACCESS_KEY}  # type: ignore

        response = requests.get(api_url, params=params).json()

        images = response["results"]
        image_url = images[random.randint(0, len(images) - 1)]["urls"]["raw"]

        image = requests.get(image_url).content

        file_system = FileSystemStorage()
        rel_path = f"{query}.jpg"
        file_system.save(rel_path, ContentFile(image))

        image_file = File(
            user=self.conversion.user,
            conversion=self.conversion,
            type=".jpg",
            file=rel_path,
            is_output=False,
        )

        image_file.save()
        return image_file

    # Executed in parallel
    def build_slide(self, slide_num: int) -> SlideContent:
        image_slide_likelihood = {0: 0, 1: 0.2, 2: 0.3, 3: 0.4, 4: 0.6, 5: 0.7, 6: 0.8}

        is_image_slide = (
            False  # random.random() < image_slide_likelihood[self.image_frequency]
        )

        if slide_num == 1:
            layout = self.manager.get_title_slide_layout()
            fields = self.manager.get_slide_layout_fields(layout)

            content: dict[SlideFieldTypes, typing.Union[str, File]] = {}

            for key, value in fields.items():
                if key == SlideFieldTypes.TITLE:
                    content[
                        key
                    ] = "This is the title"  # self.prompt_assistant(self.prompts["title"])
                elif (
                    key == SlideFieldTypes.TEXT
                ):  # self.prompt_assistant(f"{self.prompts['sub-title']}{title}")
                    content[key] = "This is the sub-title"

            return SlideContent(slide_num, SlideTypes.TITLE, layout, fields, content)

        elif is_image_slide:
            layout = self.manager.get_image_slide_layout()
            fields = self.manager.get_slide_layout_fields(layout)

            content = {}

            for key, value in fields.items():
                if key == SlideFieldTypes.TITLE:
                    content[
                        key
                    ] = "This is the title"  # self.prompt_assistant(self.prompts["title"])
                elif key == SlideFieldTypes.TEXT:
                    content[
                        key
                    ] = "This is some text"  # self.prompt_assistant(self.prompts["image-caption"])
                elif key == SlideFieldTypes.IMAGE:
                    query = (
                        "HELLO"  # self.prompt_assistant(self.prompts["image-search"])
                    )
                    content[key] = File()  # self.image_search(query)

            return SlideContent(slide_num, SlideTypes.IMAGE, layout, fields, content)

        else:
            layout = self.manager.get_content_slide_layout()
            fields = self.manager.get_slide_layout_fields(layout)

            content = {}
            for key, value in fields.items():
                if key == SlideFieldTypes.TITLE:
                    content[
                        key
                    ] = "This is the title"  # self.prompt_assistant(self.prompts["title"])
                elif key == SlideFieldTypes.TEXT:
                    content[
                        key
                    ] = "This is the text"  # self.prompt_assistant(self.prompts["content"])

            return SlideContent(slide_num, SlideTypes.CONTENT, layout, fields, content)

    def build_presentation(self) -> str:
        file_system = FileSystemStorage()

        # template = Presentation(file_system.path(f"template_{self.template}.pptx"))
        # template = Presentation(file_system.path(f"delete-my-slides.pptx"))
        self.manager.setup(file_system.path(f"template_{self.template}.pptx"))

        slide_contents: list[SlideContent] = []

        # Fetch slide content in parallel for speed
        # with ThreadPoolExecutor() as executor:
        #     future_slides = executor.map(
        #         self.build_slide, range(1, self.num_slides + 1)
        #     )
        #     slide_contents = list(future_slides)

        for i in range(1, self.num_slides + 1):
            slide_contents.append(self.build_slide(i))

        # Sort and add slides to presentation
        sorted_slide_contents = sorted(
            slide_contents, key=lambda slide: slide.slide_num
        )

        for slide_content in sorted_slide_contents:
            self.manager.add_slide_to_presentation(slide_content)

        return self.manager.save_presentation(self.conversion.id)

    # delete all files and assistants when PresentationGenerator object is deleted
    def __del__(self) -> None:
        for file in self.openai_files:
            self.delete_file(file.id)

        self.delete_assistant(self.assistant.id)

    # REQUESTS WITH EXPONENTIAL BACKOFF
    # All requests need to be done like this to deal with rate-limiting
    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def add_message_to_thread(self, thread_id: str, content: str) -> None:
        self.client.beta.threads.messages.create(
            thread_id=thread_id, role="user", content=content
        )

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def create_thread(self):  # type: ignore
        return self.client.beta.threads.create()

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def create_run(self, thread_id: str):  # type: ignore
        return self.client.beta.threads.runs.create(
            thread_id=thread_id, assistant_id=self.assistant.id
        )

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def retrieve_run(self, thread_id: str, run_id: str):  # type: ignore
        return self.client.beta.threads.runs.retrieve(
            thread_id=thread_id, run_id=run_id
        )

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def list_messages(self, thread_id: str):  # type: ignore
        return self.client.beta.threads.messages.list(thread_id=thread_id)

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def delete_thread(self, thread_id: str) -> None:
        self.client.beta.threads.delete(thread_id=thread_id)

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def list_files(self):  # type: ignore
        return self.client.files.list()

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def delete_file(self, file_id: str) -> None:
        self.client.files.delete(file_id)

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def list_assistants(self):  # type: ignore
        return self.client.beta.assistants.list(limit=100).data

    @retry(wait=wait_random_exponential(min=1, max=400), stop=stop_after_attempt(100))
    def delete_assistant(self, assistant_id: str) -> None:
        self.client.beta.assistants.delete(assistant_id)
