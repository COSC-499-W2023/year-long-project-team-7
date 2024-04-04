from enum import Enum
import typing
import json
from django.core.files.storage import FileSystemStorage
from django.core.files.base import ContentFile
from .models import File
from pptx import Presentation  # type: ignore
from io import BytesIO
from pptx.enum.shapes import PP_PLACEHOLDER  # type: ignore
from pptx.enum.shapes import MSO_SHAPE_TYPE
import random

class FieldTypes(Enum):
    TITLE = "QUESTION"
    TEXT = "ANSWER"


class ExerciseField:
    def __init__(
        self, field_index: int, field_type: FieldTypes, value: typing.Union[str, File]
    ):
        self.field_index = field_index
        self.field_type = field_type
        self.value = value


class ExerciseContent:
    def __init__(
        self,
        slide_num: int,
        layout: typing.Any,
        fields: list[ExerciseField],
    ):
        self.slide_num = slide_num
        self.layout = layout
        self.fields = fields

    def to_json_string(self) -> str:
        slide_json = {
            "SLIDE_NUM": self.slide_num,
            "SLIDE_LAYOUT": self.layout.name,
            "FIELDS": [
                {
                    "FIELD_INDEX": field.field_index,
                    "FIELD_TYPE": field.field_type,
                    "FIELD_VALUE": field.value,
                }
                for field in self.fields
            ],
        }
        return str(slide_json)

    def update_from_json(self, response: str) -> None:
        object = json.loads(response)
        for field in self.fields:
            for new_field in object["FIELDS"]:
                if str(field.field_index) == new_field["FIELD_INDEX"]:
                    field.value = new_field["FIELD_VALUE"]


class MissingPlaceholderError(Exception):
    def __init__(self, message: str):
        self.message = "Error: " + message
        super().__init__(self.message)


class ExerciseManager:
    def setup(self, template: File) -> None:
        file_system = FileSystemStorage()

        self.template_path = file_system.path(template.file.name)
        self.exercise = Presentation(self.template_path)

        self.delete_all_slides()

    def delete_all_slides(self) -> None:
        # Delete all slides from template presentation
        xml_slides = self.exercise.slides._sldIdLst
        slides = list(xml_slides)
        for slide in slides:
            rId = slide.get("r:id")
            if rId is not None:
                self.exercise.part.drop_rel(rId)
        xml_slides.clear()

    def save_presentation(self, exercise_id: int) -> str:
        file_system = FileSystemStorage()
        buffer = BytesIO()
        self.exercise.save(buffer)
        buffer.seek(0)
        file_content = ContentFile(buffer.read())

        rel_path = f"exercise_output_{exercise_id}.pptx"

        file_system.save(rel_path, file_content)

        return rel_path
    
    def get_content_slide_layout(self) -> typing.Any:
        potential_layouts = []

        temp_presentation = Presentation(self.template_path)

        for slide_layout in temp_presentation.slide_layouts:
            temp_slide = temp_presentation.slides.add_slide(slide_layout)

            for shape in temp_slide.shapes:
                if shape.is_placeholder:
                    phf = shape.placeholder_format
                    if phf.type == PP_PLACEHOLDER.BODY:
                        potential_layouts.append(slide_layout)
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    potential_layouts.append(slide_layout)
        try:
            return random.choice(potential_layouts)
        except IndexError:
            raise MissingPlaceholderError(
                "Template does not have any text placeholders"
            )

    def get_slide_layout_fields_multiple_choice(self, layout: typing.Any) -> list[ExerciseField]:
        fields = []
        temp_presentation = Presentation(self.template_path)
        temp_slide = temp_presentation.slides.add_slide(layout)
        shapes = temp_slide.shapes
        for shape in shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER.TITLE:
                    fields.append(
                        ExerciseField(
                            shapes.index(shape), FieldTypes.TITLE, "<SLIDE QUESTION HERE>"
                        )
                    )

                elif phf.type == PP_PLACEHOLDER.BODY:
                    fields.append(
                        ExerciseField(
                            shapes.index(shape), FieldTypes.TEXT, "A) <SLIDE ANSWER HERE>  B) <SLIDE ANSWER HERE>  C) <SLIDE ANSWER HERE>  D) <SLIDE ANSWER HERE>"
                        )
                    )
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                fields.append(
                    ExerciseField(
                        shapes.index(shape), FieldTypes.TEXT, "A) <SLIDE ANSWER HERE>  B) <SLIDE ANSWER HERE>  C) <SLIDE ANSWER HERE>  D) <SLIDE ANSWER HERE>"
                    )
                )

        return fields
    
    def get_slide_layout_fields_true_false(self, layout: typing.Any) -> list[ExerciseField]:
        fields = []
        temp_presentation = Presentation(self.template_path)
        temp_slide = temp_presentation.slides.add_slide(layout)
        shapes = temp_slide.shapes
        for shape in shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER.TITLE:
                    fields.append(
                        ExerciseField(
                            shapes.index(shape), FieldTypes.TITLE, "<SLIDE QUESTION HERE>"
                        )
                    )

                elif phf.type == PP_PLACEHOLDER.BODY:
                    fields.append(
                        ExerciseField(
                            shapes.index(shape), FieldTypes.TEXT, "A) TRUE  B) FALSE"
                        )
                    )
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                fields.append(
                    ExerciseField(
                        shapes.index(shape), FieldTypes.TEXT, "A) TRUE  B) FALSE"
                    )
                )

        return fields
    
    def get_slide_layout_fields_short_ans(self, layout: typing.Any) -> list[ExerciseField]:
        fields = []
        temp_presentation = Presentation(self.template_path)
        temp_slide = temp_presentation.slides.add_slide(layout)
        shapes = temp_slide.shapes
        for shape in shapes:
            if shape.is_placeholder:
                phf = shape.placeholder_format
                if phf.type == PP_PLACEHOLDER.TITLE:
                    fields.append(
                        ExerciseField(
                            shapes.index(shape), FieldTypes.TITLE, "<SLIDE QUESTION HERE>"
                        )
                    )

                elif phf.type == PP_PLACEHOLDER.BODY:
                    fields.append(
                        ExerciseField(
                            shapes.index(shape), FieldTypes.TEXT, "<SLIDE ANSWER HERE>"
                        )
                    )
                continue

            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                fields.append(
                    ExerciseField(
                        shapes.index(shape), FieldTypes.TEXT, "<SLIDE ANSWER HERE>"
                    )
                )

        return fields

    def add_slide_to_presentation(self, slide_content: ExerciseContent) -> None:
        slide = self.exercise.slides.add_slide(slide_content.layout)
        fields = slide_content.fields

        for field in fields:
            if field.field_type == FieldTypes.TITLE:
                slide.shapes[field.field_index].text = field.value

            elif field.field_type == FieldTypes.TEXT:
                slide.shapes[field.field_index].text = field.value